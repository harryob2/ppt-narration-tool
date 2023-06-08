from flask import Flask, url_for, render_template, request, send_file, session
import os
from pptx import Presentation
from pptx.util import Inches
import collections
import collections.abc
import shutil
import gunicorn
import tempfile

app = Flask(__name__)
app.config['SECRET_KEY'] = 'thekey'
app.config['SESSION_TYPE'] = 'filesystem'
app.config['UPLOAD_FOLDER'] = 'tmp'


@app.route('/', methods=['GET', 'POST'])
def index():
    temp_dir = create_temp_dir()
    session['temp_dir'] = temp_dir

    audio_path = os.path.join(temp_dir, 'audio')
    os.mkdir(audio_path)
    session['audio_path'] = audio_path
    session.modified = True

    return render_template('pptx_upload.html')



# Create a new temporary directory for each user's files
def create_temp_dir():
    temp_dir = tempfile.mkdtemp()
    return temp_dir

def empty(folder):                                  # this function emptys the folder
    if os.path.exists(folder):                      # first it checks if the folder exists
        for item in os.listdir(folder):             # then it iterates through every item in the folder
            item_path = os.path.join(folder, item)  # gets the path
            if os.path.isfile(item_path):           # checks if the item is a file
                os.remove(item_path)                # removes file

            elif os.path.isdir(item_path):          # else if folder
                shutil.rmtree(item_path)            # delete folder

# Update the make_pptx() function
@app.route('/make_pptx', methods=['GET', 'POST'])
def make_pptx():
    if request.method == 'POST':
        files = request.files.getlist('file')  # get list of all uploaded files

        # Create a new temporary directory
        temp_dir = session.get('temp_dir')
        audio_path = session.get('audio_path')  

        for file in files:  # iterate through every file
            filename = file.filename
            if filename == '':
                return 'No file selected'

            # save audio files in audio folder
            elif filename.endswith('.mp3') or filename.endswith('.m4a') or filename.endswith('.wav'):
                path = os.path.join(audio_path, filename)
                file.save(path)
                print(f'Audio file saved: {path}')

            # save pptx file
            else:
                pptx_path = os.path.join(temp_dir, filename)
                session['pptx_path'] = pptx_path
                session.modified = True
                file.save(pptx_path)
                print(f'Powerpoint saved: {pptx_path}')

    return 'Files uploaded successfully'
   
        
@app.route('/process_data', methods=['GET', 'POST'])
def process_data():
    temp_dir = session.get('temp_dir')
    pptx_path = session.get('pptx_path')
    audio_path = session.get('audio_path')

    narrated_file = make_narrated_pptx(pptx_path, audio_path, temp_dir)
    session['narrated_file'] = narrated_file
    return render_template('download.html', filename=narrated_file)


@app.route('/download')
def download():
    filename = request.args.get('filename')
    print('Successful download')
    return send_file(filename, as_attachment=True)




def make_narrated_pptx(pptx_path, audio_folder_path, temp_dir):
    print(f'Path for pptx passed to make_narrated_pptx(): {pptx_path}')
    print(f'Path for audio folder passed as: {audio_folder_path}')
    print(f'Temporary directory passed as: {temp_dir}')
    left = top = width = height = Inches(0.2)
    picture_path = r"static/mic.png"
    prs = Presentation(pptx_path)
    audio_folder = fr"{audio_folder_path}"

    for filename in os.listdir(audio_folder): # go through every audio file
        audio_path = os.path.join(audio_folder, filename) # create path to audio file
        file_number = int(os.path.splitext(filename)[0]) - 1 # get slide number from audio file name
        prs.slides[file_number].shapes.add_movie(audio_path, # add audio file to slide
                                                    left,top,width,height,
                                                    poster_frame_image=picture_path,
                                                    mime_type='video/unknown')

    narrated_path = os.path.join(temp_dir, 'narrated.pptx')  # Use temporary directory to store the narrated file
    print((f'Final narrated path saved as: {narrated_path}'))
    prs.save(narrated_path)  # save file to path
    return narrated_path
    

  
if __name__ == '__main__':
    app.run(debug = True)
