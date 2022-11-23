from flask import Flask, redirect, url_for, render_template, request, flash, send_file
import os
import io
from werkzeug.utils import secure_filename
from flask_wtf import FlaskForm
from flask_wtf.file import FileAllowed
from wtforms import FileField, SubmitField, MultipleFileField
from wtforms.validators import InputRequired
import collections
import collections.abc
from pptx import Presentation
from pptx.opc.package import PartFactory
from pptx.parts.media import MediaPart
from pptx.util import Inches

external_stylesheets = ["https://codepen.io/chriddyp/pen/bWLwgP.css"]

app = Flask(__name__)
app.config['SECRET_KEY'] = 'thekey'
app.config['UPLOAD_FOLDER'] = 'tmp'



class UploadPPTXFile(FlaskForm):
    file = FileField('File', validators=[InputRequired(), FileAllowed(['pptx'])])
    submit = SubmitField('Upload File')
    
class UploadAudioFile(FlaskForm):
    file = MultipleFileField('Files Upload', validators=[InputRequired(), FileAllowed(['mp3', 'm4a'])])
    submit = SubmitField('Upload File')


    
@app.route('/', methods=['GET','POST'])
@app.route('/pptx_upload', methods=['GET', 'POST'])
def pptx_upload():
    form = UploadPPTXFile()
    if form.validate_on_submit():
        file = form.file.data # First grab the file
        global path_to_pptx
        path_to_pptx = os.path.join(os.path.abspath(os.path.dirname(__file__)),app.config['UPLOAD_FOLDER'],secure_filename(file.filename))
        file.save(path_to_pptx) # Then save the file
        return redirect(url_for('audio_upload'))
    return render_template('pptx_upload.html', form=form)

@app.route('/audio_upload', methods=['GET', 'POST'])
def audio_upload():
    form = UploadAudioFile()
    if form.validate_on_submit():
        global path_to_pptx
        path_to_audio_folder = os.path.join(os.path.abspath(os.path.dirname(__file__)),app.config['UPLOAD_FOLDER'],'audio_folder')
        os.mkdir(path_to_audio_folder)
        for file in form.file.data:
            file_filename = secure_filename(file.filename)
            path_to_audio = os.path.join(path_to_audio_folder, file_filename)
            file.save(path_to_audio)
        make_narrated_pptx(path_to_pptx, path_to_audio_folder)
        return send_file('tmp/narrated.pptx', as_attachment=True)
    return render_template('audio_upload.html', form=form)

def make_narrated_pptx(pptx_path, audio_folder_path):
    left = top = width = height = Inches(1.0)
    picture_path = r"static/mic.jpeg"
    prs = Presentation(fr"{pptx_path}")
    audio_folder = fr"{audio_folder_path}"

    for filename in os.listdir(audio_folder):
        if filename.endswith(".mp3") or filename.endswith(".wav") or filename.endswith(".m4a"):
            audio_path = os.path.join(audio_folder, filename)
            file_number = int(os.path.splitext(filename)[0]) - 1
            prs.slides[file_number].shapes.add_movie(audio_path, 
                                                     left,top,width,height,
                                                     poster_frame_image=picture_path,
                                                     mime_type='video/unknown')
    prs.save('tmp/narrated.pptx') 

























#ALLOWED_EXTENSIONS = {'pptx'}
#UPLOAD_FOLDER = io.

#@app.route('/', methods=['GET', 'POST'])
#def upload_file():
#    if request.method == 'POST':
#        # check if the post request has the file part
#        if 'file' not in request.files:
#            flash('No file part')
#            return redirect(request.url)
#        file = request.files['file']
#        # If the user does not select a file, the browser submits an
#        # empty file without a filename.
#        if file.filename == '':
#            flash('No selected file')
#            return redirect(request.url)
#        if file and allowed_file(file.filename):
#            filename = secure_filename(file.filename)
#            file.save(os.path.join(app.config['UPLOAD_FOLDER'], filename))
#            return redirect(url_for('download_file', name=filename))
#    return '''
#    <!doctype html>
#    <title>Upload new File</title>
#    <h1>Upload new File</h1>
#    <form method=post enctype=multipart/form-data>
#      <input type=file name=file>
#      <input type=submit value=Upload>
#    </form>
#    '''    

#@app.route('/uploader', methods = ['GET', 'POST'])
#def upload_file():
#   if request.method == 'POST':
#      f = request.files['file']
#      f.save(secure_filename(f.filename))
#      return 'file uploaded successfully'

  




#def allowed_file(filename):
#return '.' in filename and \
#        filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS


  
if __name__ == '__main__':
    app.run(debug = True)