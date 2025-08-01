# Using older Pillow version that may have different libjpeg dependencies

from flask import Flask, url_for, render_template, request, send_file, session, jsonify
import os
from pptx import Presentation
from pptx.util import Inches
import collections
import collections.abc
import io
import base64
import tempfile
import zipfile
from werkzeug.utils import secure_filename

app = Flask(__name__)
app.config['SECRET_KEY'] = 'thekey'
app.config['MAX_CONTENT_LENGTH'] = 50 * 1024 * 1024  # 50MB max file size

# In-memory storage for serverless environment
file_storage = {}

@app.route('/', methods=['GET', 'POST'])
def index():
    return render_template('pptx_upload.html')

@app.route('/make_pptx', methods=['GET', 'POST'])
def make_pptx():
    if request.method == 'POST':
        try:
            files = request.files.getlist('file')
            
            if not files or all(file.filename == '' for file in files):
                return jsonify({'error': 'No files selected'}), 400
            
            # Generate a unique session ID
            session_id = request.form.get('session_id', 'default')
            if session_id not in file_storage:
                file_storage[session_id] = {'pptx': None, 'audio_files': {}}
            
            for file in files:
                filename = secure_filename(file.filename)
                if filename == '':
                    continue
                
                # Read file content into memory
                file_content = file.read()
                
                if filename.endswith(('.mp3', '.m4a', '.wav')):
                    # Store audio file in memory
                    file_storage[session_id]['audio_files'][filename] = file_content
                    print(f'Audio file stored in memory: {filename}')
                else:
                    # Store PPTX file in memory
                    file_storage[session_id]['pptx'] = file_content
                    file_storage[session_id]['pptx_filename'] = filename
                    print(f'PowerPoint stored in memory: {filename}')
            
            return jsonify({'success': True, 'message': 'Files uploaded successfully'})
            
        except Exception as e:
            print(f"Error in make_pptx: {str(e)}")
            return jsonify({'error': str(e)}), 500
    
    return jsonify({'error': 'Method not allowed'}), 405

@app.route('/process_data', methods=['GET', 'POST'])
def process_data():
    try:
        session_id = request.form.get('session_id', 'default')
        
        if session_id not in file_storage or not file_storage[session_id]['pptx']:
            return jsonify({'error': 'No PowerPoint file found'}), 400
        
        # Create narrated PowerPoint in memory
        narrated_file_data = make_narrated_pptx_in_memory(
            file_storage[session_id]['pptx'],
            file_storage[session_id]['audio_files']
        )
        
        # Store the narrated file
        file_storage[session_id]['narrated_pptx'] = narrated_file_data
        
        return render_template('download.html', session_id=session_id)
        
    except Exception as e:
        print(f"Error in process_data: {str(e)}")
        return jsonify({'error': str(e)}), 500

@app.route('/download')
def download():
    try:
        session_id = request.args.get('session_id', 'default')
        
        if session_id not in file_storage or 'narrated_pptx' not in file_storage[session_id]:
            return jsonify({'error': 'No narrated file found'}), 404
        
        narrated_data = file_storage[session_id]['narrated_pptx']
        
        # Create a file-like object from the data
        file_stream = io.BytesIO(narrated_data)
        file_stream.seek(0)
        
        return send_file(
            file_stream,
            mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation',
            as_attachment=True,
            download_name='narrated.pptx'
        )
        
    except Exception as e:
        print(f"Error in download: {str(e)}")
        return jsonify({'error': str(e)}), 500

def make_narrated_pptx_in_memory(pptx_data, audio_files):
    """Create narrated PowerPoint in memory without file system operations"""
    try:
        # Create a temporary file-like object for the PPTX
        pptx_stream = io.BytesIO(pptx_data)
        pptx_stream.seek(0)
        
        # Load the presentation
        prs = Presentation(pptx_stream)
        
        # Get the mic.png from static folder
        mic_path = os.path.join(os.path.dirname(__file__), 'static', 'mic.png')
        if not os.path.exists(mic_path):
            mic_path = None  # Use default poster frame if mic.png not found
        
        left = top = width = height = Inches(0.2)
        
        # Process audio files
        for filename, audio_data in audio_files.items():
            try:
                # Extract slide number from filename (assuming format like "1.mp3", "2.wav", etc.)
                slide_number = int(filename.split('.')[0]) - 1
                
                if 0 <= slide_number < len(prs.slides):
                    # Create a temporary file for the audio
                    with tempfile.NamedTemporaryFile(delete=False, suffix=os.path.splitext(filename)[1]) as temp_audio:
                        temp_audio.write(audio_data)
                        temp_audio.flush()
                        
                        # Add audio to slide
                        prs.slides[slide_number].shapes.add_movie(
                            temp_audio.name,
                            left, top, width, height,
                            poster_frame_image=mic_path,
                            mime_type='video/unknown'
                        )
                        
                        # Clean up temp file
                        os.unlink(temp_audio.name)
                        
            except (ValueError, IndexError) as e:
                print(f"Error processing audio file {filename}: {str(e)}")
                continue
        
        # Save to memory
        output_stream = io.BytesIO()
        prs.save(output_stream)
        output_stream.seek(0)
        
        return output_stream.getvalue()
        
    except Exception as e:
        print(f"Error in make_narrated_pptx_in_memory: {str(e)}")
        raise

@app.route('/health')
def health():
    """Health check endpoint for Vercel"""
    return jsonify({'status': 'healthy'}), 200

if __name__ == '__main__':
    app.run(debug=True) 